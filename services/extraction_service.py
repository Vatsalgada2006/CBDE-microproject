import os
import mimetypes
from typing import Optional, Tuple
import logging

logger = logging.getLogger(__name__)

# OCR threshold: if extracted text length is less than this, consider OCR for PDF/images
OCR_THRESHOLD = 50

class ExtractionService:
    def __init__(self):
        # Initialize any necessary resources
        pass

    def extract_text(self, file_path: str, content_type: str = None, filename: str = None) -> Tuple[str, str]:
        """
        Extract text from a file.
        Returns a tuple (extracted_text, status) where status is one of:
        'success', 'unsupported', 'failed'
        """
        if content_type is None and filename:
            content_type, _ = mimetypes.guess_type(filename)
            if content_type is None:
                content_type = 'application/octet-stream'

        # First, try standard extraction
        standard_text, standard_status = self._try_standard_extraction(file_path, content_type, filename)

        # Determine if we should attempt OCR
        should_try_ocr = False
        if standard_status == 'unsupported':
            # Try OCR for unsupported types that might be images or PDFs
            if self._is_ocr_applicable(content_type, filename):
                should_try_ocr = True
        elif standard_status == 'success' and len(standard_text.strip()) < OCR_THRESHOLD:
            # Text extracted but too little - might be a scanned PDF/image PDF
            if self._is_ocr_applicable(content_type, filename):
                should_try_ocr = True

        if should_try_ocr:
            ocr_text, ocr_status = self._extract_ocr(file_path, content_type, filename)
            if ocr_status == 'success' and len(ocr_text.strip()) >= len(standard_text.strip()):
                # OCR gave us equal or more text, use it
                logger.info(f"Using OCR extraction for {filename} (standard length: {len(standard_text)}, OCR length: {len(ocr_text)})")
                return ocr_text, 'success'
            else:
                logger.info(f"OCR did not improve extraction for {filename}; using standard result")

        # Return standard result (could be success, unsupported, or failed)
        return standard_text, standard_status

    def _try_standard_extraction(self, file_path: str, content_type: str, filename: str) -> Tuple[str, str]:
        """Attempt standard extraction based on content type/filename."""
        try:
            if content_type == 'application/pdf' or (filename and filename.lower().endswith('.pdf')):
                return self._extract_pdf(file_path), 'success'
            elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or \
                 (filename and filename.lower().endswith('.docx')):
                return self._extract_docx(file_path), 'success'
            elif content_type == 'text/plain' or (filename and filename.lower().endswith('.txt')):
                return self._extract_txt(file_path), 'success'
            elif content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or \
                 (filename and filename.lower().endswith('.pptx')):
                return self._extract_pptx(file_path), 'success'
            elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or \
                 (filename and filename.lower().endswith('.xlsx')):
                return self._extract_xlsx(file_path), 'success'
            else:
                logger.warning(f"Unsupported content type: {content_type} for file {filename}")
                return '', 'unsupported'
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {e}")
            return '', 'failed'

    def _is_ocr_applicable(self, content_type: str, filename: str) -> bool:
        """Check if OCR should be attempted for this content type/filename."""
        if content_type is None and filename:
            content_type, _ = mimetypes.guess_type(filename)
        if content_type is None:
            return False
        # OCR applicable for PDFs and image types
        if content_type == 'application/pdf':
            return True
        if content_type.startswith('image/'):
            return True
        # Also check filename extensions as fallback
        if filename:
            ext = filename.lower().split('.')[-1] if '.' in filename else ''
            if ext in ('pdf', 'png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif'):
                return True
        return False

    def _extract_ocr(self, file_path: str, content_type: str, filename: str) -> Tuple[str, str]:
        """Extract text using OCR (pytesseract)."""
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            logger.warning("OCR dependencies (pytesseract, Pillow) not installed. Skipping OCR.")
            return '', 'failed'

        try:
            # Handle PDFs: convert each page to image then OCR
            if content_type == 'application/pdf' or (filename and filename.lower().endswith('.pdf')):
                try:
                    from pdf2image import convert_from_path
                except ImportError:
                    logger.warning("pdf2image not installed. Skipping OCR for PDF.")
                    return '', 'failed'

                images = convert_from_path(file_path, dpi=300)
                text = ""
                for img in images:
                    text += pytesseract.image_to_string(img) + "\n"
                return text, 'success' if text.strip() else 'failed'

            # Handle images
            elif content_type.startswith('image/'):
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                return text, 'success' if text.strip() else 'failed'

            # Fallback: try to open as image anyway
            else:
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                return text, 'success' if text.strip() else 'failed'

        except Exception as e:
            logger.error(f"OCR error for {filename}: {e}")
            return '', 'failed'

    def _extract_pdf(self, file_path: str) -> str:
        import PyPDF2
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _extract_docx(self, file_path: str) -> str:
        import docx
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    def _extract_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def _extract_pptx(self, file_path: str) -> str:
        import pptx
        presentation = pptx.Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_xlsx(self, file_path: str) -> str:
        import openpyxl
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        text += str(cell) + " "
            text += "\n"
        return text